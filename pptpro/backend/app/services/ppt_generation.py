"""
PPT 생성 서비스 - python-pptx를 사용한 .pptx 파일 생성
"""
import io
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from app.db.memory_store import Project, Slide


class PPTTemplateRenderer:
    """PPT 템플릿별 렌더링 클래스"""
    
    def __init__(self, presentation: Presentation):
        self.prs = presentation
        # 기본 색상 팔레트
        self.colors = {
            'primary': RGBColor(0, 123, 255),      # Blue
            'secondary': RGBColor(108, 117, 125),   # Gray
            'success': RGBColor(40, 167, 69),       # Green
            'warning': RGBColor(255, 193, 7),       # Yellow
            'danger': RGBColor(220, 53, 69),        # Red
            'dark': RGBColor(52, 58, 64),           # Dark Gray
            'light': RGBColor(248, 249, 250)       # Light Gray
        }
    
    def render_message_only(self, slide, content: Dict[str, Any]):
        """메시지 중심 템플릿"""
        # 제목 추가
        title = slide.shapes.title
        title.text = content.get('main_message', '')
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 지원 포인트들
        supporting_points = content.get('supporting_points', [])
        if supporting_points:
            # 텍스트 박스 추가
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(8)
            height = Inches(4)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.clear()
            
            for i, point in enumerate(supporting_points):
                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(20)
                p.font.color.rgb = self.colors['secondary']
                p.space_after = Pt(12)
        
        # Call to Action
        cta = content.get('call_to_action', '')
        if cta:
            left = Inches(1)
            top = Inches(7)
            width = Inches(8)
            height = Inches(1)
            
            cta_box = slide.shapes.add_textbox(left, top, width, height)
            cta_frame = cta_box.text_frame
            cta_frame.text = cta
            cta_frame.paragraphs[0].font.size = Pt(18)
            cta_frame.paragraphs[0].font.color.rgb = self.colors['primary']
            cta_frame.paragraphs[0].font.bold = True
            cta_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def render_asis_tobe(self, slide, content: Dict[str, Any]):
        """As-Is To-Be 템플릿"""
        # 제목
        title = slide.shapes.title
        title.text = "As-Is vs To-Be"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        
        # As-Is 섹션 (왼쪽)
        as_is_title = content.get('as_is_title', 'As-Is')
        as_is_points = content.get('as_is_points', [])
        
        self._add_two_column_content(
            slide, 
            as_is_title, as_is_points, 
            content.get('to_be_title', 'To-Be'), content.get('to_be_points', []),
            left_color=self.colors['danger'], 
            right_color=self.colors['success']
        )
        
        # 전환 방법
        transition = content.get('transition_method', '')
        if transition:
            left = Inches(3)
            top = Inches(6.5)
            width = Inches(4)
            height = Inches(1)
            
            trans_box = slide.shapes.add_textbox(left, top, width, height)
            trans_frame = trans_box.text_frame
            trans_frame.text = f"→ {transition}"
            trans_frame.paragraphs[0].font.size = Pt(16)
            trans_frame.paragraphs[0].font.color.rgb = self.colors['primary']
            trans_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def render_case_box(self, slide, content: Dict[str, Any]):
        """케이스 박스 템플릿"""
        title = slide.shapes.title
        title.text = "Cases & Options"
        
        cases = content.get('cases', [])
        if not cases:
            return
        
        # 케이스들을 그리드로 배치
        cols = 2 if len(cases) > 2 else len(cases)
        rows = (len(cases) + cols - 1) // cols
        
        box_width = Inches(4)
        box_height = Inches(2.5)
        start_left = Inches(0.5)
        start_top = Inches(2)
        
        for i, case in enumerate(cases[:4]):  # 최대 4개까지
            col = i % cols
            row = i // cols
            
            left = start_left + col * (box_width + Inches(0.5))
            top = start_top + row * (box_height + Inches(0.3))
            
            # 케이스 박스 추가
            case_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, box_width, box_height
            )
            case_box.fill.solid()
            case_box.fill.fore_color.rgb = self.colors['light']
            case_box.line.color.rgb = self.colors['primary']
            case_box.line.width = Pt(2)
            
            # 케이스 텍스트 추가
            text_frame = case_box.text_frame
            text_frame.clear()
            
            # 제목
            p = text_frame.paragraphs[0]
            p.text = case.get('title', f'Case {i+1}')
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.colors['dark']
            
            # 설명
            desc_p = text_frame.add_paragraph()
            desc_p.text = case.get('description', '')
            desc_p.font.size = Pt(12)
            desc_p.font.color.rgb = self.colors['secondary']
    
    def render_step_flow(self, slide, content: Dict[str, Any]):
        """단계별 플로우 템플릿"""
        title = slide.shapes.title
        title.text = "Implementation Steps"
        
        steps = content.get('steps', [])
        if not steps:
            return
        
        # 단계별 화살표 플로우
        step_width = Inches(1.5)
        step_height = Inches(1.2)
        arrow_width = Inches(0.8)
        
        total_width = len(steps) * step_width + (len(steps) - 1) * arrow_width
        start_left = (Inches(10) - total_width) / 2
        top = Inches(3)
        
        for i, step in enumerate(steps):
            left = start_left + i * (step_width + arrow_width)
            
            # 단계 박스
            step_box = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left, top, step_width, step_height
            )
            step_box.fill.solid()
            step_box.fill.fore_color.rgb = self.colors['primary']
            step_box.line.color.rgb = self.colors['dark']
            
            # 단계 번호
            text_frame = step_box.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = str(step.get('order', i + 1))
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 단계 제목 (아래쪽)
            title_left = left - Inches(0.5)
            title_top = top + step_height + Inches(0.2)
            title_width = step_width + Inches(1)
            title_height = Inches(0.8)
            
            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            title_frame.text = step.get('title', f'Step {i+1}')
            title_frame.paragraphs[0].font.size = Pt(12)
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 화살표 (마지막 단계 제외)
            if i < len(steps) - 1:
                arrow_left = left + step_width
                arrow_top = top + step_height / 2
                arrow_box = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top, arrow_width, Inches(0.4)
                )
                arrow_box.fill.solid()
                arrow_box.fill.fore_color.rgb = self.colors['secondary']
    
    def render_chart_insight(self, slide, content: Dict[str, Any]):
        """차트 & 인사이트 템플릿"""
        title = slide.shapes.title
        title.text = content.get('chart_title', 'Data Insights')
        
        # 차트 영역 (왼쪽)
        chart_left = Inches(0.5)
        chart_top = Inches(2)
        chart_width = Inches(5)
        chart_height = Inches(4)
        
        chart_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, chart_left, chart_top, chart_width, chart_height
        )
        chart_placeholder.fill.solid()
        chart_placeholder.fill.fore_color.rgb = self.colors['light']
        chart_placeholder.line.color.rgb = self.colors['secondary']
        
        # 차트 플레이스홀더 텍스트
        chart_frame = chart_placeholder.text_frame
        chart_frame.text = f"[{content.get('chart_type', 'Chart')} 차트 영역]\n\n데이터 소스:\n{content.get('data_source', 'USER_NEEDED')}"
        chart_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        chart_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 인사이트 영역 (오른쪽)
        insights = content.get('key_insights', [])
        if insights:
            insight_left = Inches(6)
            insight_top = Inches(2)
            insight_width = Inches(4)
            insight_height = Inches(4)
            
            insight_box = slide.shapes.add_textbox(insight_left, insight_top, insight_width, insight_height)
            insight_frame = insight_box.text_frame
            insight_frame.clear()
            
            # 인사이트 제목
            title_p = insight_frame.paragraphs[0]
            title_p.text = "📈 Key Insights"
            title_p.font.size = Pt(18)
            title_p.font.bold = True
            title_p.font.color.rgb = self.colors['primary']
            
            # 인사이트 항목들
            for insight in insights:
                p = insight_frame.add_paragraph()
                p.text = f"• {insight}"
                p.font.size = Pt(14)
                p.font.color.rgb = self.colors['dark']
                p.space_after = Pt(8)
    
    def render_node_map(self, slide, content: Dict[str, Any]):
        """노드 맵 템플릿"""
        title = slide.shapes.title
        title.text = content.get('central_concept', 'Concept Map')
        
        # 중심 노드
        center_left = Inches(4)
        center_top = Inches(3.5)
        center_width = Inches(2)
        center_height = Inches(1)
        
        center_node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, center_left, center_top, center_width, center_height
        )
        center_node.fill.solid()
        center_node.fill.fore_color.rgb = self.colors['primary']
        center_node.line.color.rgb = self.colors['dark']
        
        center_frame = center_node.text_frame
        center_frame.text = content.get('central_concept', 'Central')
        center_frame.paragraphs[0].font.size = Pt(14)
        center_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        center_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        center_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 주변 노드들
        primary_nodes = content.get('primary_nodes', [])
        if primary_nodes:
            angles = [0, 60, 120, 180, 240, 300]  # 6개 노드까지 지원
            radius = Inches(2)
            
            for i, node_text in enumerate(primary_nodes[:6]):
                angle = angles[i] * 3.14159 / 180  # 라디안 변환
                
                node_left = center_left + center_width/2 + radius * 1.2 * (1 if angle < 3.14159/2 or angle > 3*3.14159/2 else -1) - Inches(0.75)
                node_top = center_top + center_height/2 + radius * 0.8 * (1 if angle > 0 and angle < 3.14159 else -1) - Inches(0.4)
                
                node = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, node_left, node_top, Inches(1.5), Inches(0.8)
                )
                node.fill.solid()
                node.fill.fore_color.rgb = self.colors['success']
                node.line.color.rgb = self.colors['dark']
                
                node_frame = node.text_frame
                node_frame.text = node_text
                node_frame.paragraphs[0].font.size = Pt(10)
                node_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                node_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                node_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    def _add_two_column_content(self, slide, left_title, left_points, right_title, right_points, left_color, right_color):
        """두 컬럼 콘텐츠 추가"""
        # 왼쪽 컬럼
        left_box_left = Inches(0.5)
        left_box_top = Inches(2)
        left_box_width = Inches(4)
        left_box_height = Inches(4)
        
        left_box = slide.shapes.add_textbox(left_box_left, left_box_top, left_box_width, left_box_height)
        left_frame = left_box.text_frame
        left_frame.clear()
        
        # 왼쪽 제목
        left_title_p = left_frame.paragraphs[0]
        left_title_p.text = left_title
        left_title_p.font.size = Pt(20)
        left_title_p.font.bold = True
        left_title_p.font.color.rgb = left_color
        
        # 왼쪽 포인트들
        for point in left_points:
            p = left_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['secondary']
        
        # 오른쪽 컬럼
        right_box_left = Inches(5.5)
        right_box_top = Inches(2)
        right_box_width = Inches(4)
        right_box_height = Inches(4)
        
        right_box = slide.shapes.add_textbox(right_box_left, right_box_top, right_box_width, right_box_height)
        right_frame = right_box.text_frame
        right_frame.clear()
        
        # 오른쪽 제목
        right_title_p = right_frame.paragraphs[0]
        right_title_p.text = right_title
        right_title_p.font.size = Pt(20)
        right_title_p.font.bold = True
        right_title_p.font.color.rgb = right_color
        
        # 오른쪽 포인트들
        for point in right_points:
            p = right_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors['secondary']


class PPTGenerationService:
    """PPT 생성 메인 서비스"""
    
    def __init__(self):
        self.template_renderers = {
            'message_only': PPTTemplateRenderer.render_message_only,
            'asis_tobe': PPTTemplateRenderer.render_asis_tobe,
            'case_box': PPTTemplateRenderer.render_case_box,
            'step_flow': PPTTemplateRenderer.render_step_flow,
            'chart_insight': PPTTemplateRenderer.render_chart_insight,
            'node_map': PPTTemplateRenderer.render_node_map
        }
    
    def generate_ppt(self, project: Project, slides: List[Slide]) -> io.BytesIO:
        """프로젝트와 슬라이드들로부터 PPT 생성"""
        
        # 새 프레젠테이션 생성
        prs = Presentation()
        
        # 제목 슬라이드 추가
        self._add_title_slide(prs, project)
        
        # 각 슬라이드 추가
        renderer = PPTTemplateRenderer(prs)
        
        for slide_data in sorted(slides, key=lambda x: x.order):
            if slide_data.content:  # 콘텐츠가 있는 슬라이드만 추가
                self._add_content_slide(prs, slide_data, renderer)
        
        # 마무리 슬라이드 추가
        self._add_closing_slide(prs, project)
        
        # 메모리 버퍼에 저장
        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        return ppt_buffer
    
    def _add_title_slide(self, prs: Presentation, project: Project):
        """제목 슬라이드 추가"""
        title_slide_layout = prs.slide_layouts[0]  # 제목 슬라이드 레이아웃
        slide = prs.slides.add_slide(title_slide_layout)
        
        # 제목 설정
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = project.title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        
        # 부제목에 프로젝트 정보 추가
        subtitle_text = ""
        if project.topic:
            subtitle_text += f"주제: {project.topic}\n"
        if project.target_audience:
            subtitle_text += f"대상: {project.target_audience}\n"
        if project.goal:
            subtitle_text += f"목표: {project.goal}"
        
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    
    def _add_content_slide(self, prs: Presentation, slide_data: Slide, renderer: PPTTemplateRenderer):
        """콘텐츠 슬라이드 추가"""
        # 콘텐츠 슬라이드 레이아웃 사용
        slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
        slide = prs.slides.add_slide(slide_layout)
        
        # 템플릿별 렌더링
        render_method = self.template_renderers.get(slide_data.template_type)
        if render_method:
            render_method(renderer, slide, slide_data.content)
        else:
            # 기본 렌더링
            renderer.render_message_only(slide, {
                'main_message': slide_data.head_message,
                'supporting_points': ['콘텐츠를 확인해주세요']
            })
    
    def _add_closing_slide(self, prs: Presentation, project: Project):
        """마무리 슬라이드 추가"""
        closing_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(closing_layout)
        
        # 감사 메시지
        left = Inches(2)
        top = Inches(3)
        width = Inches(6)
        height = Inches(2)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        text_frame.text = "감사합니다"
        text_frame.paragraphs[0].font.size = Pt(48)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 부가 정보
        info_top = Inches(5.5)
        info_textbox = slide.shapes.add_textbox(left, info_top, width, Inches(1))
        info_frame = info_textbox.text_frame
        info_frame.text = f"Generated by PPT Pro • {project.title}"
        info_frame.paragraphs[0].font.size = Pt(14)
        info_frame.paragraphs[0].alignment = PP_ALIGN.CENTER